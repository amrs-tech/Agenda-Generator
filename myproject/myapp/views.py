from django.shortcuts import render
from pptx import Presentation
import os
from shutil import rmtree
from django.template import RequestContext
from django.http import HttpResponseRedirect
from django.urls import reverse

from myapp.models import Document
from myapp.forms import DocumentForm


def list(request):
    # Handle file upload

    if request.method == 'POST':
        form = DocumentForm(request.POST, request.FILES)
        if form.is_valid():
            newdoc = Document(docfile = request.FILES['docfile'])
            newdoc.save()

            # Redirect to the document list after POST
            return HttpResponseRedirect(reverse('list'))
    else:
        form = DocumentForm() # A empty, unbound form

    # Load documents for the list page
    documents = Document.objects.all()


    # Render list page with the documents and the form
    return render(request, 'list.html', {'documents': documents, 'form': form})


def agenda(request):
    path = '' #Please insert the path of media directory in your app. Exmaple : /Users/aaaa/Downloads/your_app/media/documents/ - documents/ is necessary unless if you change in models
    try:
        files = os.listdir(path)
    except:
        print('Page doesn\'t exist')
        return HttpResponseRedirect('/myapp/')
    temp = []
    for i in files:
        if '.pptx' in i:
            temp.append(path+i)
    prs = Presentation(temp[-1])

    text_runs = []
    tit = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not slide.shapes.title.text == run.text:
                        text_runs.append(run.text)
                    else:
                        tit.append(slide.shapes.title.text)
    l = 0
    for i in tit:
       #count seconds for each titles
        l+=5

    for i in text_runs:
       #count seconds for each points in each slide
        l+=5
    try:
        rmtree(path)
        #print('deleted')
    except:
        print('cannot delete')
    return render(request, 'agenda.html',{'titles': tit, 'text': text_runs, 'time': l})